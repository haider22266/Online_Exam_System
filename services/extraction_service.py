from pathlib import Path
import os
import shutil

import docx
from flask import current_app
from pptx import Presentation


class ExtractionService:
    def extract(self, file_path):
        extension = Path(file_path).suffix.lower()
        if extension == ".pdf":
            return self._extract_pdf(file_path)
        if extension == ".docx":
            return self._extract_docx(file_path)
        if extension == ".pptx":
            return self._extract_pptx(file_path)
        raise ValueError("Unsupported file type.")

    def _extract_pdf(self, file_path):
        try:
            import fitz
        except ImportError as exc:
            raise ValueError("PyMuPDF is not installed correctly; PDF extraction is unavailable.") from exc

        pages = []
        try:
            with fitz.open(file_path) as pdf:
                for index, page in enumerate(pdf, start=1):
                    text = page.get_text("text")
                    if self._needs_ocr(text):
                        text = self._ocr_pdf_page(page, index)
                    pages.append({"page_number": index, "text": text})
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError("Invalid or unreadable PDF file.") from exc
        return pages

    @staticmethod
    def _needs_ocr(text):
        meaningful = "".join(character for character in (text or "") if character.isalnum())
        return len(meaningful) < current_app.config["PDF_MIN_TEXT_CHARACTERS"]

    @staticmethod
    def _ocr_pdf_page(page, page_number):
        if not current_app.config["PDF_OCR_ENABLED"]:
            return ""

        try:
            import fitz
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise ValueError(
                "This PDF is scanned and requires OCR. Install the Python OCR dependencies "
                "from requirements.txt."
            ) from exc

        command = current_app.config.get("TESSERACT_CMD", "").strip()
        if command:
            pytesseract.pytesseract.tesseract_cmd = command
        elif not shutil.which("tesseract"):
            raise ValueError(
                "This PDF is scanned and requires Tesseract OCR. Install Tesseract with the "
                "Bengali (ben) language pack, then upload the PDF again."
            )

        dpi = current_app.config["PDF_OCR_DPI"]
        scale = dpi / 72
        pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        data_dir = current_app.config.get("TESSERACT_DATA_DIR", "").strip()
        if data_dir:
            os.environ["TESSDATA_PREFIX"] = data_dir
        try:
            return pytesseract.image_to_string(
                image,
                lang=current_app.config["PDF_OCR_LANGUAGES"],
            )
        except pytesseract.TesseractError as exc:
            raise ValueError(
                f"OCR failed on PDF page {page_number}. Verify that the Bengali (ben) and "
                "English (eng) Tesseract language packs are installed."
            ) from exc

    def _extract_docx(self, file_path):
        document = docx.Document(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        return [{"page_number": 1, "text": text}]

    def _extract_pptx(self, file_path):
        presentation = Presentation(file_path)
        pages = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
            pages.append({"page_number": index, "text": "\n".join(lines)})
        return pages
